from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "deliverables" / "presentations" / "슈핏케어_프로젝트_발표자료.pptx"


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(22)


def add_two_column_slide(prs: Presentation, title: str, left: list[str], right: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(5.9), Inches(4.8))
    right_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.5), Inches(5.9), Inches(4.8))

    for box, lines, head in ((left_box, left, "현재"), (right_box, right, "개선 후")):
        tf = box.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = head
        p0.font.bold = True
        p0.font.size = Pt(28)
        for line in lines:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(20)


def main() -> None:
    prs = Presentation()

    today = datetime.now().strftime("%Y-%m-%d")
    add_title_slide(
        prs,
        "슈핏케어 챗봇 프로젝트 발표",
        f"하이브리드 맞춤 추천 챗봇 구축 결과\n발표일: {today}",
    )

    add_bullets_slide(
        prs,
        "프로젝트 목표",
        [
            "고객 이탈이 적은 대화형 발볼/핏 진단 경험 제공",
            "룰 기반 안정성과 LLM 유연성을 결합한 하이브리드 추천",
            "추천 결과를 고객이 이해하기 쉬운 문장으로 전달",
            "네이버 톡톡 채널에서 실서비스 가능한 배포 완성",
        ],
    )

    add_bullets_slide(
        prs,
        "슈핏케어 기획 의도",
        [
            "온라인 신발 구매에서 가장 큰 실패 요인인 '핏/사이즈 불확실성' 해결",
            "고객의 발 상태와 착화 경험을 대화형으로 수집해 맞춤 추천 제공",
            "상담사 경험 의존을 줄이고, 일관된 추천 품질을 시스템으로 표준화",
            "고객은 쉽게 선택하고, 운영자는 데이터로 개선하는 구조 구축",
        ],
    )

    add_bullets_slide(
        prs,
        "핵심 문제 정의",
        [
            "순차 설문형 UX로 인한 대화 이탈",
            "숫자형 버튼/전문 용어로 인한 이해도 저하",
            "추천 근거 부족으로 신뢰도 저하",
            "채널 연동 과정에서 웹훅/응답 경로 불안정",
        ],
    )

    add_bullets_slide(
        prs,
        "쇼핑몰 문제 해결 및 가치",
        [
            "문제 1) 사이즈 실패/반품 증가  ->  해결) 발볼·핏 중심 추천으로 미스매치 감소",
            "문제 2) 고객 문의 응대 부담  ->  해결) 챗봇 자동 진단으로 반복 상담 효율화",
            "문제 3) 구매 결정 어려움  ->  해결) TOP3 제안 + 쉬운 추천 이유로 선택 지원",
            "가치) 전환율 개선, 반품율 관리, 상담 비용 절감, 고객 만족도 향상",
        ],
    )

    add_two_column_slide(
        prs,
        "UX 개선 전/후",
        left=[
            "Q1, Q2 중심 설문 톤",
            "숫자 입력 위주 버튼",
            "관리자용 문구 고객 노출",
            "상품 추천 이유가 기술 용어 중심",
        ],
        right=[
            "대화형 문구 + 친절한 안내 톤",
            "텍스트형 버튼으로 통일",
            "고객 메시지와 관리자 정보 분리",
            "모바일 친화 3줄 추천 이유 + 간결한 TOP3",
        ],
    )

    add_bullets_slide(
        prs,
        "시스템 아키텍처",
        [
            "FastAPI + Uvicorn 기반 API 서버",
            "ConversationController 상태머신으로 대화 흐름 제어",
            "HybridProductRecommender: 후보 필터링 + 랭킹 + 설명",
            "CSV 데이터 소스(products/specs/tags/feedback) 연동",
            "Naver webhook + send API 이중 경로 대응",
        ],
    )

    add_bullets_slide(
        prs,
        "추천 로직 (실무형 하이브리드)",
        [
            "1단계: 태깅/메타데이터 기반 후보 추리기",
            "2단계: 진단 결과 기반 점수 랭킹(핏/사이즈/증상/가공성/인기도)",
            "3단계: 고객 친화 설명문 생성",
            "결과: TOP 3 + 모델코드 축약 표기 + 모바일 가독성 최적화",
        ],
    )

    add_bullets_slide(
        prs,
        "슈핏케어 핵심 로직 (요약)",
        [
            "입력 수집: 발 형태(넓음/무지외반/앞코), 사이즈, 착화 경험을 대화형으로 받음",
            "진단 엔진: 룰 기반으로 추천 핏/사이즈/보정 필요 여부를 먼저 계산",
            "후보 필터: 디자인과 핏 기준으로 상품 후보를 좁힘",
            "랭킹: 핏 적합도 + 사이즈 근접도 + 증상 대응 + 운영 지표를 합산해 TOP3 선정",
            "설명 생성: 고객이 이해하기 쉬운 문장으로 추천 이유를 요약해 전달",
        ],
    )

    add_bullets_slide(
        prs,
        "핵심 로직 간략 설명",
        [
            "한 줄 요약: 고객 발 상태를 먼저 해석한 뒤, '맞는 상품 3개'를 근거와 함께 제시하는 구조",
            "왜 효과적인가: 단순 인기순이 아니라 고객 상태를 우선 반영해 실패 구매를 줄임",
            "운영 장점: 규칙 기반이라 안정적이고, 데이터가 쌓일수록 가중치 튜닝으로 성능 개선 가능",
            "고객 경험: 어려운 용어를 줄이고 짧은 대화형 안내로 완주율과 이해도를 높임",
        ],
    )

    add_bullets_slide(
        prs,
        "데이터 구성 및 품질",
        [
            "실데이터 35개 상품 입력 완료",
            "CSV 정규화 스크립트로 오탈자/혼합값 자동 보정",
            "상태: products / product_fit_specs / product_tags / feedback 연동",
            "환경변수 기반 추천 가중치 조정 가능",
        ],
    )

    add_bullets_slide(
        prs,
        "배포 및 채널 연동 결과",
        [
            "Render 배포 완료 (운영 URL 연동)",
            "Naver webhook URL 및 send event 설정 완료",
            "운영 점검 API로 이벤트 유입 확인(/ops/naver-events-summary)",
            "최종 스모크 테스트: webhook 수신 + 응답 반환 정상",
        ],
    )

    add_bullets_slide(
        prs,
        "데모 시나리오",
        [
            "시작: 상품 먼저 추천받기 / 발 정보 입력 후 추천받기 선택",
            "진단: 발 형태/사이즈/착화 경험 입력",
            "결과: 추천 핏·사이즈·설명 + TOP3 상품 제시",
            "재진단: 같은 채팅창에서 1/2 재입력으로 재시작",
        ],
    )

    add_bullets_slide(
        prs,
        "기대 효과",
        [
            "고객 이해도 향상: 어려운 용어 제거 + 친절한 문구",
            "상담 효율 향상: 진단 자동화로 반복 문의 감소",
            "추천 신뢰도 향상: 근거 기반 결과 제공",
            "운영 확장성 확보: 데이터 누적 기반 고도화 가능",
        ],
    )

    add_bullets_slide(
        prs,
        "향후 고도화 로드맵",
        [
            "배송/교환반품 FAQ의 단계적 RAG 도입",
            "주문/반품 폐루프 데이터와 추천 검증 로직 확장",
            "A/B 테스트로 문구/버튼/가중치 지속 최적화",
            "운영 대시보드 고도화(이벤트/전환/반품율 추적)",
        ],
    )

    add_bullets_slide(
        prs,
        "SaaS 확장 방향",
        [
            "멀티 쇼핑몰/브랜드 지원: shop_id 기반 정책/카탈로그 분리",
            "도메인 확장: 여성화 -> 남성화/아동화/특정 카테고리로 확장",
            "구독형 수익모델: 기본형(진단) + 고급형(분석/대시보드/RAG)",
            "데이터 네트워크 효과: 누적 피드백 기반 추천 성능 고도화",
            "운영 제품화: 관리자 콘솔, 리포트 자동화, A/B 테스트 내장",
        ],
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Created: {OUT_PATH}")


if __name__ == "__main__":
    main()
