from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_10 = ROOT / "deliverables" / "presentations" / "슈핏케어_프로젝트_발표자료_10분.pptx"
OUT_5 = ROOT / "deliverables" / "presentations" / "슈핏케어_프로젝트_발표자료_5분.pptx"


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str], font_size: int = 22) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(font_size)


def set_slide_notes(slide, note_text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = note_text


def build_10min() -> None:
    prs = Presentation()
    today = datetime.now().strftime("%Y-%m-%d")
    add_title_slide(prs, "슈핏케어 챗봇 프로젝트 발표 (10분)", f"중년여성 타깃 맞춤 추천 챗봇\n발표일: {today}")
    set_slide_notes(prs.slides[-1], "안녕하세요. 슈핏케어 챗봇 프로젝트를 발표하겠습니다. 이번 프로젝트는 중년여성 고객의 발볼/핏 불편을 줄이는 데 집중했습니다.")

    add_bullets_slide(prs, "프로젝트 목표", [
        "중년여성 고객의 발볼/핏 불편을 대화형 진단으로 해결",
        "룰 기반 안정성과 하이브리드 추천 결합",
        "고객 친화 문구로 이해도와 완주율 개선",
        "네이버 톡톡 실서비스 배포 완료",
    ])
    set_slide_notes(prs.slides[-1], "프로젝트 목표는 대화형 진단 경험, 하이브리드 추천, 이해 쉬운 결과 전달, 그리고 네이버 채널 실서비스 완성입니다.")
    add_bullets_slide(prs, "기획 의도 및 문제 정의", [
        "온라인 신발 구매의 핵심 문제: 사이즈 불확실성과 반품",
        "중년여성 고객군은 무지외반/발볼 압박 이슈가 빈번",
        "설문형 UX와 어려운 용어로 이탈 증가",
        "상담 품질을 시스템으로 표준화할 필요",
    ])
    set_slide_notes(prs.slides[-1], "중년여성 고객은 발볼 압박, 무지외반 등으로 온라인 구매 실패가 잦습니다. 이 문제를 대화형 진단으로 해결하는 것이 기획 의도입니다.")
    add_bullets_slide(prs, "해결 전략", [
        "대화형 입력: 텍스트 버튼 중심으로 진입 장벽 완화",
        "하이브리드 추천: 후보 필터링 + 점수 랭킹 + 설명 생성",
        "고객용 메시지 최적화: 관리자 문구 제거, 쉬운 설명 제공",
        "재진단 지원: 같은 채팅창에서 즉시 재시작",
    ])
    set_slide_notes(prs.slides[-1], "해결 전략은 UX 개선, 하이브리드 추천, 고객용 문구 개선, 재진단 지원 네 가지입니다.")
    add_bullets_slide(prs, "핵심 로직", [
        "1) 입력 수집: 발 형태/사이즈/착화 경험",
        "2) 진단 엔진: 추천 핏·사이즈·보정 계산",
        "3) 후보 필터링: 디자인/핏 기준으로 후보 축소",
        "4) TOP3 랭킹: 적합도 중심 점수 합산",
        "5) 설명 생성: 모바일 친화 3줄 요약 + 간결한 상품 목록",
    ])
    set_slide_notes(prs.slides[-1], "핵심 로직은 입력 수집에서 시작해 진단, 후보 필터링, TOP3 랭킹, 쉬운 설명 전달까지 이어지는 파이프라인입니다.")
    add_bullets_slide(prs, "데이터 구성", [
        "실데이터 35개 상품 입력 완료",
        "CSV 4종 연동: 상품/핏스펙/태그/피드백",
        "정규화 스크립트로 오탈자·혼합값 자동 보정",
        "가중치 환경변수화로 운영 중 튜닝 가능",
    ])
    set_slide_notes(prs.slides[-1], "실데이터 35개를 연동했고, CSV 정규화 자동화로 데이터 품질을 보완했습니다. 운영 중 가중치 조정도 가능합니다.")
    add_bullets_slide(prs, "아키텍처 및 연동", [
        "FastAPI + Uvicorn + 상태머신(ConversationController)",
        "HybridProductRecommender 모듈 분리",
        "Naver webhook + send API 경로 동시 지원",
        "운영 점검 API(/ops/naver-events-summary)로 실시간 확인",
    ])
    set_slide_notes(prs.slides[-1], "FastAPI 기반으로 상태머신과 추천기를 분리했고, 네이버 webhook과 send API를 모두 지원하도록 구성했습니다.")
    add_bullets_slide(prs, "배포/검증 결과", [
        "Render 배포 완료 및 최신 커밋 반영 확인",
        "네이버 웹훅 수신/응답/이벤트 증가 확인",
        "최종 스모크 테스트 1회 통과",
        "실서비스 가능한 안정 상태 도달",
    ])
    set_slide_notes(prs.slides[-1], "Render 배포와 네이버 연동을 완료했고, 운영 점검 API로 이벤트 유입과 응답 정상 동작을 확인했습니다.")
    add_bullets_slide(prs, "기대 효과", [
        "고객: 이해 쉬운 안내 + 실패 구매 감소",
        "쇼핑몰: 상담 효율 향상 + 반품율 관리",
        "운영: 데이터 기반 지속 개선 가능",
        "비즈니스: 전환율 개선 및 재구매 가능성 확대",
    ])
    set_slide_notes(prs.slides[-1], "고객은 이해하기 쉬운 추천을 받고, 쇼핑몰은 상담 효율과 반품 관리 측면에서 효과를 기대할 수 있습니다.")
    add_bullets_slide(prs, "SaaS 확장 방향", [
        "멀티 쇼핑몰/브랜드 지원 (shop_id 기반 분리)",
        "카테고리 확장 (여성화 -> 타 카테고리)",
        "구독형 상품화 (기본형/고급형 분석)",
        "폐루프 데이터 기반 추천 고도화",
    ])
    set_slide_notes(prs.slides[-1], "향후에는 멀티 쇼핑몰 지원, 카테고리 확장, 구독형 상품화, 데이터 기반 성능 고도화를 추진합니다.")
    add_bullets_slide(prs, "결론", [
        "슈핏케어는 중년여성 타깃의 실사용 가능한 맞춤 추천 챗봇",
        "대화형 진단 + 하이브리드 추천으로 고객 경험과 운영 효율 동시 개선",
        "현재 목표 범위 100% 완료, 고도화 단계 진입 가능",
    ])
    set_slide_notes(prs.slides[-1], "정리하면 슈핏케어는 중년여성 타깃 문제를 실서비스 가능한 형태로 해결한 프로젝트이며, 현재 범위 목표를 완료했습니다.")

    prs.save(str(OUT_10))


def build_5min() -> None:
    prs = Presentation()
    today = datetime.now().strftime("%Y-%m-%d")
    add_title_slide(prs, "슈핏케어 챗봇 프로젝트 발표 (5분)", f"중년여성 타깃 맞춤 추천 핵심 요약\n발표일: {today}")
    set_slide_notes(prs.slides[-1], "안녕하세요. 5분 안에 슈핏케어 프로젝트의 핵심만 요약해 말씀드리겠습니다.")

    add_bullets_slide(prs, "문제와 목표", [
        "문제: 중년여성 온라인 신발 구매 시 발볼/핏 불확실성으로 반품 증가",
        "목표: 대화형 진단으로 맞는 핏·사이즈 추천",
        "결과: 이해 쉬운 설명과 TOP3 제안으로 선택 부담 완화",
    ], font_size=24)
    set_slide_notes(prs.slides[-1], "중년여성 고객의 온라인 구매 실패 문제를 줄이고, 발 상태 기반 추천으로 선택 부담을 낮추는 것이 목표입니다.")
    add_bullets_slide(prs, "핵심 해결 방식", [
        "대화형 UX: 쉬운 버튼 중심 입력",
        "하이브리드 추천: 후보 필터 + 점수 랭킹 + 설명 생성",
        "고객 친화 출력: 3줄 추천 이유 + 간결한 TOP3",
    ], font_size=24)
    set_slide_notes(prs.slides[-1], "핵심은 대화형 UX와 하이브리드 추천입니다. 고객이 쉽게 입력하고 이해할 수 있도록 메시지를 단순화했습니다.")
    add_bullets_slide(prs, "핵심 로직 한눈에", [
        "입력 수집 -> 진단 엔진 -> 후보 필터링 -> TOP3 랭킹 -> 설명 전달",
        "단순 인기순이 아닌 고객 발 상태 우선 추천",
        "재진단 가능 구조로 반복 개선",
    ], font_size=24)
    set_slide_notes(prs.slides[-1], "입력 수집부터 TOP3 제안까지 하나의 흐름으로 동작하며, 단순 인기순이 아닌 고객 상태를 우선 반영합니다.")
    add_bullets_slide(prs, "구현/배포 성과", [
        "실데이터 35개 연동 + 정규화 자동화",
        "Render 배포 및 네이버 톡톡 연동 완료",
        "웹훅 수신/응답 최종 스모크 테스트 통과",
    ], font_size=24)
    set_slide_notes(prs.slides[-1], "실데이터 연동, Render 배포, 네이버 연동, 최종 스모크 테스트까지 완료해 실사용 가능한 상태입니다.")
    add_bullets_slide(prs, "기대 효과 및 확장", [
        "고객: 사이즈 실패 감소, 이해도 향상",
        "쇼핑몰: 상담 효율화, 반품 관리",
        "확장: SaaS 멀티 테넌트/FAQ-RAG/폐루프 고도화",
    ], font_size=24)
    set_slide_notes(prs.slides[-1], "고객 경험과 운영 효율을 동시에 개선하고, 이후 SaaS 구조로 확장 가능한 기반을 마련했습니다.")
    add_bullets_slide(prs, "마무리", [
        "슈핏케어는 중년여성 타깃 문제를 실제 운영 가능한 형태로 해결",
        "현재 범위 구축 완료, 다음 단계는 데이터 기반 고도화",
        "감사합니다.",
    ], font_size=24)
    set_slide_notes(prs.slides[-1], "이상으로 발표를 마치겠습니다. 감사합니다.")

    prs.save(str(OUT_5))


def main() -> None:
    OUT_10.parent.mkdir(parents=True, exist_ok=True)
    build_10min()
    build_5min()
    print(f"Created: {OUT_10}")
    print(f"Created: {OUT_5}")


if __name__ == "__main__":
    main()
